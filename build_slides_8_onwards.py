"""
Builds slides 8-14 for the ABC Benchmarking MAS thesis defence.
Matches the dark professional theme of the existing Canva deck.
Run:  python build_slides_8_onwards.py
Output: slides_8_onwards.pptx  (import into Canva to replace slides 8-end)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os, copy
from pptx.oxml.ns import qn
from lxml import etree

FIGURES = os.path.join(os.path.dirname(__file__), "figures")

# ── palette ────────────────────────────────────────────────────────────────
BG          = RGBColor(0x0F, 0x19, 0x2B)   # very dark navy
BG_CARD     = RGBColor(0x1A, 0x2A, 0x40)   # card surface
BG_CARD2    = RGBColor(0x0D, 0x22, 0x38)   # alternate card
COL_TITLE   = RGBColor(0xFF, 0xFF, 0xFF)
COL_BODY    = RGBColor(0xCB, 0xD5, 0xE1)
COL_MUTED   = RGBColor(0x94, 0xA3, 0xB8)
COL_BLUE    = RGBColor(0x3B, 0x82, 0xF6)   # condition A
COL_ORANGE  = RGBColor(0xF9, 0x73, 0x16)   # condition B
COL_PURPLE  = RGBColor(0xA7, 0x8B, 0xFA)   # condition C
COL_GREEN   = RGBColor(0x22, 0xC5, 0x5E)   # positive delta
COL_ACCENT  = RGBColor(0x38, 0xBD, 0xF8)   # highlight

W  = Inches(13.33)   # LAYOUT_WIDE
H  = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

def blank():
    layout = prs.slide_layouts[6]   # completely blank
    sl = prs.slides.add_slide(layout)
    bg = sl.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return sl

def txb(sl, text, x, y, w, h, size=16, bold=False, color=None,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tf = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf.word_wrap = wrap
    p = tf.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or COL_BODY
    return tf

def title_block(sl, number, title, subtitle=None):
    """Slide number badge + big title + optional subtitle."""
    # Number badge
    badge = sl.shapes.add_shape(1, Inches(0.45), Inches(0.30), Inches(0.52), Inches(0.52))
    badge.fill.solid(); badge.fill.fore_color.rgb = COL_BLUE
    badge.line.fill.background()
    tf = badge.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(number)
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = COL_TITLE

    # Title
    txb(sl, title, 1.10, 0.22, 11.5, 0.72,
        size=30, bold=True, color=COL_TITLE, align=PP_ALIGN.LEFT)
    if subtitle:
        txb(sl, subtitle, 1.10, 0.88, 11.5, 0.38,
            size=13, color=COL_MUTED, align=PP_ALIGN.LEFT)

def card(sl, x, y, w, h, color=None):
    shp = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = color or BG_CARD
    shp.line.fill.background()
    return shp

def hline(sl, y):
    ln = sl.shapes.add_shape(
        9, Inches(0.45), Inches(y), Inches(12.43), Inches(0.01))
    ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0x2D, 0x3F, 0x55)
    ln.line.fill.background()

def img(sl, fname, x, y, w, h):
    path = os.path.join(FIGURES, fname)
    if os.path.exists(path):
        sl.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))

def stat_card(sl, x, y, w, h, value, label, val_color=None):
    card(sl, x, y, w, h, BG_CARD)
    txb(sl, value, x+0.12, y+0.12, w-0.24, h*0.52,
        size=34, bold=True, color=val_color or COL_ACCENT,
        align=PP_ALIGN.CENTER)
    txb(sl, label, x+0.08, y + h*0.58, w-0.16, h*0.38,
        size=11, color=COL_MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 8 — Updated Experimental Design
# ══════════════════════════════════════════════════════════════════════
sl8 = blank()
title_block(sl8, "08", "Experimental Design",
            "Five controlled sweeps — isolating threshold, judge, and architecture effects")
hline(sl8, 1.15)

# Row labels
cols = [
    ("sweep_phase1",     "Phase 1 Baseline",       "A / B / C",  "0.70",  "qwen3-32b / qwen3-32b"),
    ("sweep_thresh85_b", "Threshold Sensitivity B", "B only",     "0.85",  "qwen3-32b / qwen3-32b"),
    ("sweep_thresh85_c", "Threshold Sensitivity C", "C only",     "0.85",  "qwen3-32b / qwen3-32b"),
    ("sweep_judge_sep_a","Judge Separation A",       "A only",     "—",     "qwen2.5-7b / qwen3-32b"),
    ("sweep_judge_sep_b","Judge Separation B",       "B only",     "0.85",  "qwen2.5-7b / qwen3-32b"),
]

headers = ["Sweep File", "Purpose", "Conditions", "Threshold", "Generator / Judge"]
col_x   = [0.45, 2.55, 6.05, 7.45, 8.65]
col_w   = [2.00, 3.40, 1.30, 1.10, 4.30]

# Header row
card(sl8, 0.45, 1.22, 12.43, 0.52, RGBColor(0x1E, 0x3A, 0x5F))
for i, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    txb(sl8, hdr, cx+0.08, 1.27, cw-0.12, 0.42,
        size=11, bold=True, color=COL_ACCENT, align=PP_ALIGN.LEFT)

row_colors = [BG_CARD, BG_CARD2, BG_CARD, BG_CARD2, BG_CARD]
for ri, (sweep, purpose, cond, thresh, models) in enumerate(cols):
    ry = 1.74 + ri * 0.88
    card(sl8, 0.45, ry, 12.43, 0.82, row_colors[ri])
    # file tag
    tag = sl8.shapes.add_shape(1, Inches(col_x[0]+0.08), Inches(ry+0.18),
                                Inches(1.70), Inches(0.45))
    tag.fill.solid(); tag.fill.fore_color.rgb = RGBColor(0x0D, 0x3B, 0x6E)
    tag.line.fill.background()
    tf2 = tag.text_frame; tf2.word_wrap = False
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = sweep
    r2.font.size = Pt(9); r2.font.bold = True
    r2.font.color.rgb = RGBColor(0x7D, 0xD3, 0xFC)

    for i, (val, cx, cw) in enumerate(zip(
            [" ", purpose, cond, thresh, models], col_x, col_w)):
        if i == 0:
            continue
        align = PP_ALIGN.CENTER if i in (2, 3) else PP_ALIGN.LEFT
        c = COL_ORANGE if (i == 3 and thresh == "0.85") else \
            COL_GREEN  if (i == 4 and "7b" in models)  else COL_BODY
        txb(sl8, val, cx+0.08, ry+0.22, cw-0.12, 0.50,
            size=12, color=c, align=align)

txb(sl8, "All sweeps: 19 SWE-bench Lite instances across django, sympy, scikit-learn, astropy",
    0.45, 6.70, 12.43, 0.40, size=11, color=COL_MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 9 — Phase 1 Headline Numbers
# ══════════════════════════════════════════════════════════════════════
sl9 = blank()
title_block(sl9, "09", "Phase 1 Results — Headline Numbers",
            "19 instances × 3 conditions — qwen3-32b generator and judge")
hline(sl9, 1.15)

# 5 stat cards top row
stats = [
    ("0.878", "A  mean score\n(σ = 0.145, min 0.30)", COL_BLUE),
    ("0.881", "B @ 0.70  mean\n(σ = 0.069, min 0.72)", COL_ORANGE),
    ("0.926", "B @ 0.85  mean\n(σ = 0.044, min 0.85)", COL_GREEN),
    ("0.887", "C @ 0.70  mean\n(σ = 0.126, min 0.40)", COL_PURPLE),
    ("0.721", "C @ 0.85  mean\n(σ = 0.353, min 0.00)", RGBColor(0xF8, 0x71, 0x71)),
]
card_w = 2.38
for i, (val, lbl, col) in enumerate(stats):
    stat_card(sl9, 0.45 + i*(card_w+0.08), 1.30, card_w, 1.55, val, lbl, col)

# Key findings below cards
findings = [
    (COL_GREEN,  "▲ B@0.85 beats B@0.70",
     "Mean +0.045, stdev halved (0.069→0.044), hard floor at 0.85. "
     "Raising the threshold forces more refine iterations and eliminates all below-threshold answers."),
    (COL_ORANGE, "✓ Reviewer narrows variance",
     "A spans 0.30–0.96; B@0.85 spans 0.85–1.00. The refine loop is effectively a "
     "score floor: no instance escapes without satisfying the rubric."),
    (COL_PURPLE, "⚠ Condition C is bimodal at 0.85",
     "14 / 19 instances score 0.86–1.00. The remaining 5 regress sharply (≤ 0.40), "
     "including one at 0.00. Aggressive refinement with an action agent is not uniform improvement."),
    (COL_ACCENT, "★ Judge separation amplifies lift",
     "With a weaker Ollama 7B generator the Reviewer adds Δ = +0.114 vs only Δ = +0.003 "
     "when the generator is strong — the Reviewer matters most where it is needed most."),
]
for i, (col, hdr, body) in enumerate(findings):
    fy = 3.05 + i * 1.05
    dot = sl9.shapes.add_shape(1, Inches(0.45), Inches(fy+0.14),
                                Inches(0.10), Inches(0.36))
    dot.fill.solid(); dot.fill.fore_color.rgb = col
    dot.line.fill.background()
    txb(sl9, hdr,  0.65, fy+0.08, 12.0, 0.32, size=13, bold=True,  color=col)
    txb(sl9, body, 0.65, fy+0.42, 12.0, 0.55, size=11, color=COL_BODY)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 10 — Finding 1: Reviewer Floor Enforcement
# ══════════════════════════════════════════════════════════════════════
sl10 = blank()
title_block(sl10, "10", "Finding 1 — Reviewer Enforces a Hard Score Floor",
            "Every B@0.85 instance ≥ 0.85 · Baseline A spans 0.30 → 0.96")
hline(sl10, 1.15)

# Left: figure
img(sl10, "embed_floor.png", 0.45, 1.25, 7.60, 5.85)

# Right: callout cards
callouts = [
    (COL_GREEN,  "0.30 → 0.96",
     "sympy-16792: biggest single lift.\nBaseline scored 0.30 — Reviewer drove it to 0.96."),
    (COL_ORANGE, "σ halved",
     "A stdev = 0.145 · B@0.85 stdev = 0.044.\nRefine loop compresses the score distribution."),
    (COL_ACCENT, "Hard minimum",
     "B@0.85 minimum = 0.85 across all 19 instances.\nThe rubric + refine acts as a guaranteed floor."),
    (COL_BLUE,   "73 % of instances lifted",
     "14 of 19 instances improved under B vs A.\nOnly 5 showed marginal regression (all ≥ 0.85)."),
]
for i, (col, hdr, body) in enumerate(callouts):
    cy = 1.30 + i * 1.52
    card(sl10, 8.30, cy, 4.60, 1.35, BG_CARD)
    bar = sl10.shapes.add_shape(1, Inches(8.30), Inches(cy),
                                 Inches(0.10), Inches(1.35))
    bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background()
    txb(sl10, hdr,  8.52, cy+0.10, 4.20, 0.40, size=14, bold=True, color=col)
    txb(sl10, body, 8.52, cy+0.52, 4.20, 0.72, size=11, color=COL_BODY)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 11 — Finding 2: Condition C Bimodality
# ══════════════════════════════════════════════════════════════════════
sl11 = blank()
title_block(sl11, "11", "Finding 2 — Condition C is Bimodal at Threshold 0.85",
            "Aggressive refinement with an action agent is not uniformly beneficial")
hline(sl11, 1.15)

img(sl11, "embed_bimodal.png", 0.45, 1.25, 7.30, 5.60)

rows = [
    (COL_PURPLE, "High cluster",   "14 / 19 instances",  "Score 0.86 – 1.00\nAction agent succeeds reliably"),
    (RGBColor(0xF8,0x71,0x71), "Regression tail", "5 / 19 instances",
     "Score ≤ 0.40 (one at 0.00)\nAggressive refine degrades output"),
    (COL_ORANGE, "Root cause",     "Context explosion",
     "OpenHands accumulates ~150K tokens\nby step 30 — judge loses coherence"),
    (COL_ACCENT, "Implication",    "Threshold matters",
     "C@0.70 mean = 0.887 vs C@0.85 = 0.721\nHigher threshold amplifies tail risk"),
]
for i, (col, tag, stat, body) in enumerate(rows):
    ry = 1.30 + i * 1.50
    card(sl11, 7.95, ry, 5.00, 1.35, BG_CARD)
    tag_shp = sl11.shapes.add_shape(1, Inches(7.95), Inches(ry),
                                     Inches(0.10), Inches(1.35))
    tag_shp.fill.solid(); tag_shp.fill.fore_color.rgb = col
    tag_shp.line.fill.background()
    txb(sl11, tag,  8.17, ry+0.08, 4.60, 0.32, size=12, bold=True, color=col)
    txb(sl11, stat, 8.17, ry+0.40, 4.60, 0.28, size=13, bold=True, color=COL_TITLE)
    txb(sl11, body, 8.17, ry+0.72, 4.60, 0.52, size=11, color=COL_BODY)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 12 — Finding 3: Judge Separation
# ══════════════════════════════════════════════════════════════════════
sl12 = blank()
title_block(sl12, "12", "Finding 3 — Reviewer Adds Most Value When Generator is Weak",
            "Judge separation experiment: Ollama qwen2.5-7B generator · qwen3-32B judge")
hline(sl12, 1.15)

img(sl12, "embed_judge_sep.png", 0.45, 1.28, 7.50, 5.55)

# Right: three interpretation cards
interps = [
    (COL_GREEN,  "Δ = +0.114  (weak gen)",
     "Mean A = 0.657, mean B = 0.771.\n"
     "When the generator struggles the Reviewer's\nrefine loop delivers the largest lift."),
    (COL_MUTED,  "Δ = +0.003  (strong gen)",
     "Mean A = 0.878, mean B = 0.881.\n"
     "A capable generator leaves little room\nfor the Reviewer to improve."),
    (COL_ACCENT, "Architectural implication",
     "The Reviewer is not redundant — it is a\n"
     "safety net that scales with generator weakness.\n"
     "Deploy it whenever the generator is uncertain."),
]
for i, (col, hdr, body) in enumerate(interps):
    iy = 1.30 + i * 1.95
    card(sl12, 8.15, iy, 4.75, 1.80, BG_CARD)
    bar2 = sl12.shapes.add_shape(1, Inches(8.15), Inches(iy),
                                  Inches(0.10), Inches(1.80))
    bar2.fill.solid(); bar2.fill.fore_color.rgb = col; bar2.line.fill.background()
    txb(sl12, hdr,  8.38, iy+0.15, 4.35, 0.42, size=13, bold=True, color=col)
    txb(sl12, body, 8.38, iy+0.58, 4.35, 1.10, size=11, color=COL_BODY)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 13 — Conclusions & Future Work
# ══════════════════════════════════════════════════════════════════════
sl13 = blank()
title_block(sl13, "13", "Conclusions, Limitations & Future Work")
hline(sl13, 1.15)

# Left column — conclusions
card(sl13, 0.45, 1.25, 6.00, 5.85, BG_CARD)
txb(sl13, "CONCLUSIONS", 0.65, 1.38, 5.60, 0.40,
    size=12, bold=True, color=COL_ACCENT)

conclusions = [
    ("Rubric-based gating works",
     "B@0.85 lifts mean score by +0.048 over A and enforces a hard minimum — "
     "no instance escapes below 0.85."),
    ("Reviewer value scales with generator weakness",
     "Δ = +0.003 (strong model) vs Δ = +0.114 (7B model). "
     "The Reviewer is most valuable where the generator is most uncertain."),
    ("Hallucination Bridge generalises the rubric",
     "The same 5-dimension rubric audits OpenHands action traces — a structurally "
     "distinct hallucination class becomes detectable without fine-tuning."),
    ("Bimodality is a real risk at C@0.85",
     "Aggressive refinement with an autonomous action agent can degrade output "
     "when the judge loses coherence under large context."),
]
for i, (hdr, body) in enumerate(conclusions):
    cy = 1.88 + i * 1.25
    dot = sl13.shapes.add_shape(1, Inches(0.65), Inches(cy+0.06),
                                 Inches(0.10), Inches(0.28))
    dot.fill.solid(); dot.fill.fore_color.rgb = COL_GREEN; dot.line.fill.background()
    txb(sl13, hdr,  0.88, cy, 5.40, 0.30, size=12, bold=True, color=COL_TITLE)
    txb(sl13, body, 0.88, cy+0.32, 5.40, 0.82, size=11, color=COL_BODY)

# Right column — limitations + future
card(sl13, 6.65, 1.25, 6.23, 2.60, BG_CARD2)
txb(sl13, "LIMITATIONS", 6.85, 1.38, 5.80, 0.36,
    size=12, bold=True, color=RGBColor(0xF8, 0x71, 0x71))
lims = [
    "Absolute scores are LLM-regime sensitive — only within-regime ordering is comparable.",
    "Free-tier throttling affected Condition C; token costs per run are high (~150K at step 30).",
    "19-instance set provides directional signal but limited statistical power.",
]
for i, lim in enumerate(lims):
    txb(sl13, f"– {lim}", 6.85, 1.82 + i*0.70, 5.80, 0.62, size=11, color=COL_BODY)

card(sl13, 6.65, 4.05, 6.23, 3.05, BG_CARD2)
txb(sl13, "FUTURE WORK", 6.85, 4.18, 5.80, 0.36,
    size=12, bold=True, color=COL_ORANGE)
future = [
    "Full 50-instance SWE-bench Lite sweep for statistical significance.",
    "Formal per-action hallucination metric on OpenHands traces.",
    "Threshold calibration — per-domain adaptive thresholding.",
    "Evaluate with stronger open-source models (Qwen2.5-72B, GPT-4o).",
]
for i, fw in enumerate(future):
    txb(sl13, f"→  {fw}", 6.85, 4.60 + i*0.60, 5.80, 0.52, size=11, color=COL_BODY)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 14 — Thank You
# ══════════════════════════════════════════════════════════════════════
sl14 = blank()

# Large centred thank-you block — vertically centred in upper half
txb(sl14, "THANK YOU", 1.0, 1.20, 11.33, 1.20,
    size=56, bold=True, color=COL_TITLE, align=PP_ALIGN.CENTER)
txb(sl14, "FOR YOUR TIME", 1.0, 2.32, 11.33, 0.62,
    size=28, bold=False, color=COL_ACCENT, align=PP_ALIGN.CENTER)

hline(sl14, 3.15)

txb(sl14,
    "Syed Muhammad Shajee Raza  ·  Group J4232  ·  ITMO University",
    0.45, 3.28, 12.43, 0.38,
    size=13, color=COL_MUTED, align=PP_ALIGN.CENTER)

txb(sl14,
    "Development of a Lightweight Hallucination Detection Algorithm  "
    "for the Reviewer Agent in a Code-Generation Multi-Agent System",
    0.45, 3.70, 12.43, 0.55,
    size=12, italic=True, color=COL_MUTED, align=PP_ALIGN.CENTER)

# Summary stat row at bottom
stat_data = [
    ("19",     "SWE-bench instances"),
    ("5",      "Sweep configurations"),
    ("3",      "Architectural conditions"),
    ("+0.114", "Max Reviewer lift"),
    ("0.85",   "B@0.85 score floor"),
]
sw = (13.33 - 0.90) / len(stat_data)
for i, (val, lbl) in enumerate(stat_data):
    sx = 0.45 + i * sw
    stat_card(sl14, sx, 4.50, sw - 0.10, 1.75, val, lbl,
              COL_GREEN if "+" in val else COL_ACCENT)


# ── write ──────────────────────────────────────────────────────────────────
OUT = os.path.join(os.path.dirname(__file__), "slides_8_onwards.pptx")
prs.save(OUT)
print("Saved: " + OUT + "  (" + str(len(prs.slides)) + " slides)")
